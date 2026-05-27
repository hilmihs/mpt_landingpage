"""
Generate PROPOSAL_MAJELIS.pptx — proposal deck untuk Ketua Majelis MPT.

Target audience: non-IT religious authority. Tone: formal, accessible,
analogi-rich, dengan visualisasi diagram alur dan ringkasan.
"""

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# =========================================================
# Brand palette (matches /lib brand & landing page)
# =========================================================
PRIMARY = RGBColor(0x1F, 0x29, 0x37)   # dark navy
ACCENT = RGBColor(0xB8, 0x86, 0x0B)    # gold
ACCENT_SOFT = RGBColor(0xE5, 0xC4, 0x80)
SURFACE = RGBColor(0xF8, 0xF4, 0xE9)   # cream paper
SURFACE_DEEP = RGBColor(0xEC, 0xE3, 0xCD)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT = RGBColor(0x55, 0x5A, 0x66)
INK_MUTE = RGBColor(0x95, 0x9C, 0xA8)
LINE = RGBColor(0xE5, 0xE2, 0xD8)
SUCCESS = RGBColor(0x22, 0x86, 0x3A)
WARN = RGBColor(0xC4, 0x71, 0x10)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# =========================================================
# Helpers
# =========================================================
def setup_slide(slide, bg=SURFACE):
    """Set background fill color."""
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.shadow.inherit = False
    # send to back
    sp = bg_shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    """Add a text box with single paragraph."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return box


def add_text_multi(slide, x, y, w, h, parts, anchor=MSO_ANCHOR.TOP):
    """parts = list of (text, {size, bold, color, align, font})."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    for i, (text, opts) in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        if "space_before" in opts:
            p.space_before = Pt(opts["space_before"])
        r = p.add_run()
        r.text = text
        r.font.size = Pt(opts.get("size", 16))
        r.font.bold = opts.get("bold", False)
        r.font.color.rgb = opts.get("color", INK)
        r.font.name = opts.get("font", "Calibri")
    return box


def add_box(slide, x, y, w, h, fill=WHITE, border=LINE, border_w=1.0,
            radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(border_w)
    shp.shadow.inherit = False
    if radius:
        # Tame the default huge corner radius
        shp.adjustments[0] = 0.05
    return shp


def add_circle(slide, x, y, d, fill, border=None, border_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_arrow(slide, x, y, w, h, color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def slide_header(slide, eyebrow, title, subtitle=None):
    """Common header band on each non-cover slide."""
    # Eyebrow pill
    pill = add_box(slide, Inches(0.6), Inches(0.45), Inches(2.1), Inches(0.32),
                   fill=SURFACE_DEEP, border=LINE, border_w=0.5)
    add_text(slide, Inches(0.6), Inches(0.45), Inches(2.1), Inches(0.32),
             eyebrow, size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # Title
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12.2), Inches(0.65),
             title, size=28, bold=True, color=PRIMARY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.4),
                 subtitle, size=14, color=INK_SOFT)


def slide_footer(slide, page_num, total):
    add_text(slide, Inches(0.6), Inches(7.0), Inches(8), Inches(0.3),
             "Muhajir Project Tilawah · Proposal kepada Majelis", size=9,
             color=INK_MUTE)
    add_text(slide, Inches(11.5), Inches(7.0), Inches(1.4), Inches(0.3),
             f"{page_num} / {total}", size=9, color=INK_MUTE,
             align=PP_ALIGN.RIGHT)


# =========================================================
# Build deck
# =========================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]

TOTAL = 11  # for footer

# -------- Slide 1: Cover --------
s = prs.slides.add_slide(blank)
setup_slide(s, bg=PRIMARY)

# Subtle accent corner
add_circle(s, Inches(-2.5), Inches(-2.5), Inches(6), ACCENT_SOFT)
acc_corner = s.shapes[-1]
acc_corner.fill.transparency = 0  # full
# Note: python-pptx doesn't expose transparency well; use a darker overlay instead
acc_corner.fill.solid()
acc_corner.fill.fore_color.rgb = RGBColor(0x2A, 0x35, 0x42)
acc_corner.line.fill.background()

# Bismillah Arabic
add_text(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
         "بِسْمِ اللَّهِ الرَّحْمَٰنِ الرَّحِيمِ",
         size=32, color=ACCENT, align=PP_ALIGN.CENTER,
         font="Traditional Arabic")

# Main title
add_text_multi(s, Inches(0.6), Inches(2.2), Inches(12.2), Inches(2.2), [
    ("PROPOSAL KEPADA MAJELIS", {
        "size": 13, "bold": True, "color": ACCENT,
        "align": PP_ALIGN.CENTER, "font": "Calibri"}),
    ("Asistensi Tilawah Al-Fatihah", {
        "size": 48, "bold": True, "color": WHITE,
        "align": PP_ALIGN.CENTER, "space_before": 16}),
    ("Memuliakan Bacaan Al-Quran lewat Pendampingan", {
        "size": 22, "color": ACCENT_SOFT,
        "align": PP_ALIGN.CENTER, "space_before": 8}),
    ("yang Tetap Berakar pada Ulama", {
        "size": 22, "color": ACCENT_SOFT, "align": PP_ALIGN.CENTER}),
])

# Footer band
add_box(s, Inches(0), Inches(6.7), prs.slide_width, Inches(0.8),
        fill=RGBColor(0x16, 0x1E, 0x2B), border=PRIMARY, radius=False)
add_text_multi(s, Inches(0.6), Inches(6.85), Inches(12.2), Inches(0.5), [
    ("Muhajir Project Tilawah", {
        "size": 12, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}),
    (f"Disampaikan oleh Tim MPT · {datetime.now().strftime('%B %Y')}", {
        "size": 10, "color": INK_MUTE, "align": PP_ALIGN.CENTER}),
])

# -------- Slide 2: Mukadimah --------
s = prs.slides.add_slide(blank)
setup_slide(s)
slide_header(s, "MUKADIMAH",
             "Assalamu'alaikum, dengan hormat kepada Majelis",
             "Ringkasan singkat sebelum kami uraikan teknis-nya.")

# Body card
card = add_box(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(4.5),
               fill=WHITE, border=LINE, border_w=1.0)

add_text_multi(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(4), [
    ("Banyak saudara muslim kita merasa belum yakin dengan bacaan Al-Fatihah-nya.",
     {"size": 18, "bold": True, "color": INK}),
    ("Sementara di sisi lain, jumlah pengajar yang qualified terbatas, dan kendala "
     "jadwal serta jarak membuat banyak yang akhirnya menunda — bahkan tidak pernah "
     "memperbaiki bacaannya.",
     {"size": 16, "color": INK_SOFT, "space_before": 10}),
    ("Muhajir Project Tilawah adalah ikhtiar kami menjawab kebutuhan tersebut. "
     "Bukan untuk menggantikan peran ulama, justru sebaliknya — untuk menyaring dan "
     "mengarahkan calon murid sehingga waktu pengajar dapat dimanfaatkan untuk yang "
     "paling membutuhkan.",
     {"size": 16, "color": INK_SOFT, "space_before": 14}),
    ("Proposal ini menjelaskan bagaimana cara kerjanya, prinsip syar'i yang kami jaga, "
     "dan permohonan kami kepada Majelis atas restu dan bimbingannya.",
     {"size": 16, "color": INK_SOFT, "space_before": 14}),
])

slide_footer(s, 2, TOTAL)

# -------- Slide 3: Tantangan --------
s = prs.slides.add_slide(blank)
setup_slide(s)
slide_header(s, "LATAR BELAKANG",
             "Tantangan yang Kami Hadapi",
             "Tiga hambatan utama yang dirasakan umat saat ingin memperbaiki bacaan.")

cards = [
    ("Mayoritas Belum Mantap",
     "Banyak muslim yang sudah dewasa belum merasa yakin dengan bacaan "
     "Al-Fatihah-nya — padahal Al-Fatihah dibaca minimum 17 kali setiap hari "
     "dalam shalat fardhu.",
     "01"),
    ("Pengajar Terbatas",
     "Satu pengajar tahsin yang qualified bisa kewalahan menangani ratusan calon "
     "murid. Antrean panjang, slot mengajar terbatas oleh waktu fisik.",
     "02"),
    ("Hambatan Jadwal & Jarak",
     "Banyak calon murid akhirnya menunda atau menyerah karena tidak cocok jam "
     "kerja, tinggal jauh dari pengajar, atau tidak tahu harus mulai dari mana.",
     "03"),
]

card_w_in = 3.9
card_w = Inches(card_w_in)
card_h = Inches(3.7)
card_y = Inches(2.4)
card_gap_in = 0.25
card_gap = Inches(card_gap_in)
total_w_in = card_w_in * 3 + card_gap_in * 2
start_x_in = (13.333 - total_w_in) / 2

for i, (head, body, num) in enumerate(cards):
    x = Inches(start_x_in + (card_w_in + card_gap_in) * i)
    add_box(s, x, card_y, card_w, card_h, fill=WHITE, border=LINE)
    # Number circle
    add_circle(s, x + Inches(0.4), card_y + Inches(0.4), Inches(0.7),
               fill=PRIMARY)
    add_text(s, x + Inches(0.4), card_y + Inches(0.4), Inches(0.7), Inches(0.7),
             num, size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # Title
    add_text(s, x + Inches(0.4), card_y + Inches(1.4), card_w - Inches(0.8),
             Inches(0.7), head, size=18, bold=True, color=PRIMARY)
    # Body
    add_text(s, x + Inches(0.4), card_y + Inches(2.1), card_w - Inches(0.8),
             Inches(1.8), body, size=13, color=INK_SOFT)

slide_footer(s, 3, TOTAL)

# -------- Slide 4: Visi --------
s = prs.slides.add_slide(blank)
setup_slide(s, bg=PRIMARY)

add_text(s, Inches(0.6), Inches(1.0), Inches(12.2), Inches(0.5),
         "VISI", size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Huge quote
add_text(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.8),
         "“", size=80, color=ACCENT, align=PP_ALIGN.CENTER)

add_text_multi(s, Inches(1.5), Inches(2.6), Inches(10.3), Inches(3.0), [
    ("Setiap muslim dapat memperbaiki bacaan",
     {"size": 30, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}),
    ("Al-Fatihah-nya — dengan panduan pengajar",
     {"size": 30, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}),
    ("yang qualified, mulai dari rumah.",
     {"size": 30, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}),
])

add_text(s, Inches(0.6), Inches(6.0), Inches(12.2), Inches(0.5),
         "— Visi MPT, ditujukan kepada umat Muslim di Indonesia —",
         size=12, color=ACCENT_SOFT, align=PP_ALIGN.CENTER)

slide_footer(s, 4, TOTAL)

# -------- Slide 5: Cara Kerja (Peserta Journey) --------
s = prs.slides.add_slide(blank)
setup_slide(s)
slide_header(s, "CARA KERJA",
             "Bagaimana Peserta Menggunakannya",
             "Tiga langkah sederhana — dapat dilakukan dari rumah lewat HP.")

steps = [
    ("REKAM BACAAN",
     "Peserta merekam bacaan Al-Fatihah-nya lewat HP — hanya butuh 3 menit, "
     "tidak perlu install aplikasi apapun.",
     "1"),
    ("DAPAT RAPOT AWAL",
     "AI memeriksa rekaman dan menghasilkan rapot awal dalam 30 detik. "
     "Rapot ini menunjukkan area yang perlu diperbaiki.",
     "2"),
    ("BERTEMU PENGAJAR",
     "Peserta yang ingin lanjut, akan dipertemukan dengan pengajar via Zoom "
     "untuk pendampingan langsung dan ijazah.",
     "3"),
]

step_w_in = 3.6
step_w = Inches(step_w_in)
step_h = Inches(3.4)
step_y = Inches(2.4)
arrow_w_in = 0.5
arrow_w = Inches(arrow_w_in)
gap_in = 0.05
total_w_in = step_w_in * 3 + arrow_w_in * 2 + gap_in * 4
start_x_in = (13.333 - total_w_in) / 2

for i, (head, body, num) in enumerate(steps):
    x_in = start_x_in + (step_w_in + arrow_w_in + gap_in * 2) * i
    x = Inches(x_in)
    add_box(s, x, step_y, step_w, step_h, fill=WHITE, border=LINE)
    # Big numeral
    add_text(s, x, step_y + Inches(0.3), step_w, Inches(1.5),
             num, size=84, bold=True, color=ACCENT_SOFT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.3), step_y + Inches(1.7), step_w - Inches(0.6),
             Inches(0.5), head, size=14, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), step_y + Inches(2.2), step_w - Inches(0.6),
             Inches(1.1), body, size=12, color=INK_SOFT,
             align=PP_ALIGN.CENTER)
    if i < 2:
        ax = Inches(x_in + step_w_in + gap_in)
        add_arrow(s, ax, step_y + Inches(1.5), arrow_w, Inches(0.4), color=ACCENT)

# Bottom note
add_text(s, Inches(0.6), Inches(6.2), Inches(12.2), Inches(0.4),
         "Setelah sesi pengajar, peserta dapat lanjut ke Tahsin 4-sesi → lulus → program lanjutan HITS.",
         size=13, color=INK_SOFT, align=PP_ALIGN.CENTER)

slide_footer(s, 5, TOTAL)

# -------- Slide 6: 4 Indikator --------
s = prs.slides.add_slide(blank)
setup_slide(s)
slide_header(s, "PENILAIAN AI",
             "Empat Indikator yang Diperiksa",
             "Hanya kategori Lahn Jaliy (kesalahan jelas). Lahn Khafi tetap di tangan pengajar.")

indikators = [
    ("Harakat",
     "Tanda baca yang menentukan vokal — fathah, kasrah, dammah, sukun. "
     "AI mendeteksi kesalahan pembacaan tanda baca dasar.",
     "ا َ ِ ُ"),
    ("Huruf",
     "Pengucapan huruf yang benar — makhraj (tempat keluar suara) dan sifat dasar. "
     "AI mendeteksi tukar huruf yang sering terjadi.",
     "ع ح ق ك"),
    ("Panjang-Pendek",
     "Hukum mad — kapan suara harus dipanjangkan dan berapa lama. "
     "AI mendeteksi mad yang dipendekkan atau sebaliknya.",
     "ـا ـو ـي"),
    ("Syaddah",
     "Penekanan ganda pada huruf bertasydid. "
     "AI mendeteksi syaddah yang terlewat atau ditambahkan keliru.",
     "ـّ"),
]

grid_w_in = 5.7
grid_w = Inches(grid_w_in)
grid_h_in = 2.0
grid_h = Inches(grid_h_in)
grid_y_in = 2.3
gap_in = 0.25
start_x_in = (13.333 - grid_w_in * 2 - gap_in) / 2

for i, (head, body, glyph) in enumerate(indikators):
    col = i % 2
    row = i // 2
    x = Inches(start_x_in + (grid_w_in + gap_in) * col)
    y = Inches(grid_y_in + (grid_h_in + gap_in) * row)
    add_box(s, x, y, grid_w, grid_h, fill=WHITE, border=LINE)
    # Arabic glyph badge
    add_box(s, x + Inches(0.3), y + Inches(0.3), Inches(1.4), Inches(1.4),
            fill=SURFACE_DEEP, border=ACCENT_SOFT)
    add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(1.4), Inches(1.4),
             glyph, size=26, color=PRIMARY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font="Traditional Arabic")
    # Text
    add_text(s, x + Inches(1.85), y + Inches(0.35), grid_w - Inches(2.1),
             Inches(0.5), head, size=18, bold=True, color=PRIMARY)
    add_text(s, x + Inches(1.85), y + Inches(0.85), grid_w - Inches(2.1),
             Inches(1.1), body, size=12, color=INK_SOFT)

# Disclaimer band
add_box(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
        fill=SURFACE_DEEP, border=LINE, border_w=0.5)
add_text(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
         "AI hanya mendeteksi 4 kategori di atas. Keputusan ijazah & evaluasi mendalam tetap di tangan pengajar.",
         size=11, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

slide_footer(s, 6, TOTAL)

# -------- Slide 7: AI vs Pengajar --------
s = prs.slides.add_slide(blank)
setup_slide(s)
slide_header(s, "PEMBAGIAN PERAN",
             "AI sebagai Asisten — BUKAN Pengganti Pengajar",
             "Setiap peran punya batas yang jelas. AI menyaring, pengajar memutuskan.")

col_w_in = 5.9
col_w = Inches(col_w_in)
col_h = Inches(3.9)
col_y = Inches(2.3)
gap_in = 0.3
gap = Inches(gap_in)
start_x_in = (13.333 - col_w_in * 2 - gap_in) / 2

# AI column
ai_x = Inches(start_x_in)
add_box(s, ai_x, col_y, col_w, col_h, fill=WHITE, border=LINE)
add_box(s, ai_x, col_y, col_w, Inches(0.55), fill=PRIMARY, border=PRIMARY)
add_text(s, ai_x, col_y, col_w, Inches(0.55),
         "AI  (Mesin)", size=15, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

ai_items = [
    "Identifikasi kesalahan dasar (4 indikator)",
    "Berikan laporan tekstual dalam 30 detik",
    "Tersedia 24/7, dapat diakses kapan saja",
    "Konsisten — tidak lelah, tidak emosional",
    "Menyaring siapa yang siap untuk pengajar",
]
for i, item in enumerate(ai_items):
    y = col_y + Inches(0.85 + i * 0.55)
    add_text(s, ai_x + Inches(0.3), y, Inches(0.3), Inches(0.4),
             "•", size=18, bold=True, color=ACCENT, anchor=MSO_ANCHOR.TOP)
    add_text(s, ai_x + Inches(0.65), y, col_w - Inches(0.9), Inches(0.5),
             item, size=13, color=INK)

# Pengajar column
ust_x = Inches(start_x_in + col_w_in + gap_in)
add_box(s, ust_x, col_y, col_w, col_h, fill=WHITE, border=ACCENT, border_w=2.5)
add_box(s, ust_x, col_y, col_w, Inches(0.55), fill=ACCENT, border=ACCENT)
add_text(s, ust_x, col_y, col_w, Inches(0.55),
         "PENGAJAR  (Manusia)", size=15, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

ust_items = [
    "Talaqqi langsung — sanad turun ke murid",
    "Koreksi adab, niat, dan tarbiyah ruhi",
    "Evaluasi mendalam (Lahn Khafi & lain-lain)",
    "Memberi ijazah dan sertifikasi",
    "Keputusan akhir — di sisi syar'i & teknis",
]
for i, item in enumerate(ust_items):
    y = col_y + Inches(0.85 + i * 0.55)
    add_text(s, ust_x + Inches(0.3), y, Inches(0.3), Inches(0.4),
             "✓", size=18, bold=True, color=SUCCESS, anchor=MSO_ANCHOR.TOP)
    add_text(s, ust_x + Inches(0.65), y, col_w - Inches(0.9), Inches(0.5),
             item, size=13, color=INK)

# Bottom insight
add_box(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.55),
        fill=PRIMARY, border=PRIMARY)
add_text(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.55),
         "AI menyaring → Pengajar memutuskan. Efisien, tapi tetap berakar pada ulama.",
         size=14, bold=True, color=ACCENT_SOFT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

slide_footer(s, 7, TOTAL)

# -------- Slide 8: Funnel Bertahap --------
s = prs.slides.add_slide(blank)
setup_slide(s)
slide_header(s, "JALUR PENDAMPINGAN",
             "Penyaringan Bertahap — 5 Tahap",
             "Setiap tahap menyaring. Yang lolos ke program lanjutan hanya yang benar-benar siap.")

stages = [
    ("Rapot AI", "Gratis, 30 detik", "Semua peserta", PRIMARY),
    ("Sesi Assessment", "60 menit · 12 peserta", "Yang tertarik", PRIMARY),
    ("Tahsin Al-Fatihah", "4 sesi × 90 menit", "Yang siap lanjut", ACCENT),
    ("Lulus Tahsin", "≥ 3 dari 4 sesi", "Yang istiqomah", ACCENT),
    ("Program HITS", "Lanjutan tilawah", "Alumni qualified", SUCCESS),
]

stage_w_in = 2.25
stage_w = Inches(stage_w_in)
stage_h = Inches(3.2)
stage_y = Inches(2.4)
arrow_w_s_in = 0.18
arrow_w_s = Inches(arrow_w_s_in)
gap_s_in = 0.05
gap_s = Inches(gap_s_in)
total_w_in = stage_w_in * 5 + arrow_w_s_in * 4 + gap_s_in * 8
start_x_in = (13.333 - total_w_in) / 2

for i, (name, detail, who, color) in enumerate(stages):
    x_in = start_x_in + (stage_w_in + arrow_w_s_in + gap_s_in * 2) * i
    x = Inches(x_in)
    # Vertical card
    add_box(s, x, stage_y, stage_w, stage_h, fill=WHITE, border=LINE)
    # Color stripe at top
    add_box(s, x, stage_y, stage_w, Inches(0.4), fill=color, border=color)
    # Step number
    add_text(s, x, stage_y, stage_w, Inches(0.4),
             f"TAHAP {i + 1}", size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Name
    add_text(s, x + Inches(0.15), stage_y + Inches(0.6), stage_w - Inches(0.3),
             Inches(0.8), name, size=15, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)
    # Detail
    add_text(s, x + Inches(0.15), stage_y + Inches(1.35), stage_w - Inches(0.3),
             Inches(0.5), detail, size=11, color=INK_SOFT,
             align=PP_ALIGN.CENTER)
    # Divider
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             x + Inches(0.6), stage_y + Inches(2.0),
                             stage_w - Inches(1.2), Inches(0.02))
    div.fill.solid()
    div.fill.fore_color.rgb = LINE
    div.line.fill.background()
    # Who
    add_text(s, x + Inches(0.15), stage_y + Inches(2.15), stage_w - Inches(0.3),
             Inches(0.3), "Yang Ikut:", size=9, bold=True,
             color=INK_MUTE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), stage_y + Inches(2.45), stage_w - Inches(0.3),
             Inches(0.6), who, size=12, color=INK, align=PP_ALIGN.CENTER)

    if i < 4:
        ax = Inches(x_in + stage_w_in + gap_s_in)
        add_arrow(s, ax, stage_y + Inches(1.5), arrow_w_s, Inches(0.3),
                  color=ACCENT)

# Note
add_text(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.7),
         "HITS Linktree (program lanjutan) HANYA muncul untuk peserta yang lulus Tahsin Al-Fatihah. "
         "Tidak ada jalan pintas — gating ini dirancang menjaga kualitas alumni.",
         size=12, color=INK_SOFT, align=PP_ALIGN.CENTER)

slide_footer(s, 8, TOTAL)

# -------- Slide 9: Adab & Syar'i --------
s = prs.slides.add_slide(blank)
setup_slide(s)
slide_header(s, "ADAB & SYAR'I",
             "Prinsip yang Kami Jaga",
             "Enam aturan adab yang melekat dalam desain sistem — bukan tempelan.")

principles = [
    ("Gender-matched STRICT",
     "Ikhwan dengan pengajar ikhwan, akhwat dengan pengajar akhwat. "
     "Sistem memfilter otomatis, tidak ada celah."),
    ("Audio Peserta Dihapus 7 Hari",
     "Rekaman bacaan otomatis dihapus dari server dalam 7 hari. "
     "Tidak disimpan permanen, tidak diolah untuk training apapun."),
    ("AI Tidak Beri Ijazah",
     "AI hanya menghasilkan rapot teknis. Sertifikasi/ijazah hanya "
     "dikeluarkan oleh pengajar manusia."),
    ("Keputusan Akhir di Pengajar",
     "Setiap kasus ambigu otomatis diserahkan ke pengajar untuk override "
     "manual. Pengajar bisa membatalkan keputusan AI kapanpun."),
    ("Data Tidak Dijual",
     "Data peserta tidak dibagikan ke pihak ketiga, tidak diolah untuk "
     "iklan, tidak dimonetisasi."),
    ("Source Code Terbuka untuk Audit",
     "Source code dapat di-audit oleh tim teknis yang ditunjuk Majelis "
     "kapan saja. Tidak ada 'kotak hitam'."),
]

p_w_in = 5.95
p_w = Inches(p_w_in)
p_h_in = 1.45
p_h = Inches(p_h_in)
p_y_start_in = 2.3
p_gap_in = 0.15
start_x_in = (13.333 - p_w_in * 2 - p_gap_in) / 2

for i, (head, body) in enumerate(principles):
    col = i % 2
    row = i // 2
    x = Inches(start_x_in + (p_w_in + p_gap_in) * col)
    y = Inches(p_y_start_in + (p_h_in + p_gap_in) * row)
    add_box(s, x, y, p_w, p_h, fill=WHITE, border=LINE)
    # Checkmark badge
    add_circle(s, x + Inches(0.25), y + Inches(0.25), Inches(0.55),
               fill=SUCCESS)
    add_text(s, x + Inches(0.25), y + Inches(0.25), Inches(0.55), Inches(0.55),
             "✓", size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Text
    add_text(s, x + Inches(1.0), y + Inches(0.2), p_w - Inches(1.2), Inches(0.4),
             head, size=14, bold=True, color=PRIMARY)
    add_text(s, x + Inches(1.0), y + Inches(0.62), p_w - Inches(1.2),
             Inches(0.75), body, size=11, color=INK_SOFT)

slide_footer(s, 9, TOTAL)

# -------- Slide 10: Tim & Tata Kelola --------
s = prs.slides.add_slide(blank)
setup_slide(s)
slide_header(s, "TIM & TATA KELOLA",
             "Tiga Lapis Pengawasan",
             "Dari syar'i hingga teknis — setiap lapis ada penanggung jawab yang jelas.")

layers = [
    ("LAPIS 1 — SYAR'I",
     "Penanggung Jawab Syar'i",
     "Diharapkan ditunjuk oleh Majelis atau dari kalangan ulama yang ditunjuk "
     "Majelis. Bertanggung jawab atas adab, fiqh, dan kualitas keagamaan secara "
     "keseluruhan. Memiliki hak veto.",
     ACCENT),
    ("LAPIS 2 — PENDIDIKAN",
     "Tim Pengajar (Ustadz/Ustadzah)",
     "Para alumni sanad bersertifikasi yang menjadi penjaga kualitas talaqqi. "
     "Mereka yang berhadapan langsung dengan peserta dan memberi ijazah. "
     "Saat ini target perekrutan: 4-6 pengajar awal.",
     PRIMARY),
    ("LAPIS 3 — TEKNIS",
     "Tim Pengembang",
     "Hilmi Sobandi (project lead) dan tim teknis. Bertanggung jawab atas "
     "infrastruktur, keamanan data, dan kelangsungan operasional platform. "
     "Tidak mengambil keputusan syar'i.",
     INK_SOFT),
]

l_w = Inches(12.1)
l_h_in = 1.35
l_h = Inches(l_h_in)
l_y_start_in = 2.3
l_gap_in = 0.15

for i, (eyebrow, head, body, color) in enumerate(layers):
    y = Inches(l_y_start_in + (l_h_in + l_gap_in) * i)
    add_box(s, Inches(0.6), y, l_w, l_h, fill=WHITE, border=LINE)
    # Color stripe (left)
    add_box(s, Inches(0.6), y, Inches(0.12), l_h, fill=color, border=color,
            radius=False)
    # Eyebrow
    add_text(s, Inches(0.9), y + Inches(0.2), Inches(3), Inches(0.3),
             eyebrow, size=10, bold=True, color=color)
    # Title
    add_text(s, Inches(0.9), y + Inches(0.5), Inches(4), Inches(0.4),
             head, size=15, bold=True, color=PRIMARY)
    # Body
    add_text(s, Inches(5.0), y + Inches(0.2), l_w - Inches(4.5), Inches(1.1),
             body, size=12, color=INK_SOFT)

# Bottom note
add_text(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.3),
         "Tim teknis tidak akan mengambil keputusan syar'i. Setiap perubahan kebijakan adab → harus melalui Lapis 1.",
         size=11, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

slide_footer(s, 10, TOTAL)

# -------- Slide 11: Permohonan --------
s = prs.slides.add_slide(blank)
setup_slide(s, bg=PRIMARY)

add_text(s, Inches(0.6), Inches(0.8), Inches(12.2), Inches(0.5),
         "PERMOHONAN", size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(1.25), Inches(12.2), Inches(0.8),
         "Mohon Restu & Bimbingan Majelis",
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

requests = [
    ("01",
     "Restu untuk Meluncurkan",
     "Restu Majelis agar proyek ini dapat hadir sebagai sarana dakwah — "
     "membantu lebih banyak muslim memperbaiki bacaan Al-Fatihah-nya."),
    ("02",
     "Bimbingan dalam Pemilihan",
     "Bimbingan Majelis dalam memilih Penanggung Jawab Syar'i dan "
     "pengajar yang qualified — ini krusial untuk kualitas keseluruhan."),
    ("03",
     "Evaluasi Berkala",
     "Evaluasi setiap 3-6 bulan agar Majelis dapat memantau adab tetap "
     "terjaga dan memberi koreksi bila perlu. Kami siap akuntabel."),
]

r_w_in = 11.5
r_w = Inches(r_w_in)
r_h_in = 1.2
r_h = Inches(r_h_in)
r_y_in = 2.55
r_gap_in = 0.18
r_start_x_in = (13.333 - r_w_in) / 2
r_start_x = Inches(r_start_x_in)

for i, (num, head, body) in enumerate(requests):
    y = Inches(r_y_in + (r_h_in + r_gap_in) * i)
    # Subtle bg
    add_box(s, r_start_x, y, r_w, r_h, fill=RGBColor(0x29, 0x33, 0x42),
            border=RGBColor(0x35, 0x40, 0x52), border_w=1.0)
    # Number
    add_text(s, r_start_x + Inches(0.4), y, Inches(1.2), r_h,
             num, size=42, bold=True, color=ACCENT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # Head
    add_text(s, r_start_x + Inches(1.7), y + Inches(0.25), Inches(8), Inches(0.4),
             head, size=16, bold=True, color=WHITE)
    # Body
    add_text(s, r_start_x + Inches(1.7), y + Inches(0.65), Inches(9.5), Inches(0.5),
             body, size=11, color=ACCENT_SOFT)

# Closing
add_text(s, Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.5),
         "بَارَكَ اللَّهُ فِيكُمْ  ·  Barakallahu fiikum",
         size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
         font="Traditional Arabic")

# =========================================================
# Save
# =========================================================
out = Path(__file__).resolve().parent.parent / "docs" / "PROPOSAL_MAJELIS.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"✓ Saved {out} ({out.stat().st_size // 1024} KB · {TOTAL} slides)")
