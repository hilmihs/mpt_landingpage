"""
Generate DECK_WALKTHROUGH.pptx — perjalanan peserta & pengajar, ujung ke ujung.

Untuk tim internal. Bukan proposal: ini dokumentasi apa yang sudah berjalan di
production per 4 Agustus 2026, memakai tangkapan layar asli.

Tangkapan layar dikumpulkan lewat Playwright dari production dan disimpan di
docs/img/deck/. Tiga gambar WhatsApp adalah rekayasa TAMPILAN — kalimatnya
disalin apa adanya dari lib/whatsapp.ts.

Jalankan:  .venv/bin/python scripts/generate-deck-walkthrough.py
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# =========================================================
# Palet — sama dengan generate-proposal-deck.py dan landing page
# =========================================================
PRIMARY = RGBColor(0x1F, 0x29, 0x37)
ACCENT = RGBColor(0xB8, 0x86, 0x0B)
ACCENT_SOFT = RGBColor(0xE5, 0xC4, 0x80)
SURFACE = RGBColor(0xF8, 0xF4, 0xE9)
SURFACE_DEEP = RGBColor(0xEC, 0xE3, 0xCD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT = RGBColor(0x55, 0x5A, 0x66)
INK_MUTE = RGBColor(0x95, 0x9C, 0xA8)
LINE = RGBColor(0xE5, 0xE2, 0xD8)
SUCCESS = RGBColor(0x22, 0x86, 0x3A)
WARN = RGBColor(0xC4, 0x71, 0x10)
DANGER = RGBColor(0xB8, 0x65, 0x4A)

ROOT = Path(__file__).resolve().parent.parent
IMG = ROOT / "docs" / "img" / "deck"
OUT = ROOT / "docs" / "DECK_WALKTHROUGH.pptx"

W_IN = 13.333
H_IN = 7.5

prs = Presentation()
prs.slide_width = Inches(W_IN)
prs.slide_height = Inches(H_IN)

TOTAL = 23
_page = 0


# =========================================================
# Helper
# =========================================================
def new_slide(bg=SURFACE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    # Latar harus di belakang semua isi
    tree = slide.shapes._spTree
    tree.remove(rect._element)
    tree.insert(2, rect._element)
    return slide


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_box(slide, x, y, w, h, fill=WHITE, border=LINE, border_w=1.0, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if radius:
        shape.adjustments[0] = 0.05
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_w)
    shape.shadow.inherit = False
    return shape


def add_circle(slide, x, y, d, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
    shape.shadow.inherit = False
    return shape


def add_arrow(slide, x, y, w, h, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def header(slide, eyebrow, title, subtitle=None):
    if eyebrow:
        pill = add_box(slide, 0.6, 0.42, max(1.1, 0.115 * len(eyebrow) + 0.34), 0.3,
                       fill=SURFACE_DEEP, border=None)
        pill.adjustments[0] = 0.5
        add_text(slide, 0.6, 0.42, pill.width / 914400, 0.3, eyebrow.upper(), size=9.5,
                 bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 0.6, 0.83, 12.2, 0.5, title, size=26, bold=True, color=PRIMARY)
    if subtitle:
        add_text(slide, 0.6, 1.4, 12.2, 0.42, subtitle, size=12.5, color=INK_SOFT)


def footer(slide, dark=False):
    global _page
    _page += 1
    c = ACCENT_SOFT if dark else INK_MUTE
    add_text(slide, 0.6, 6.95, 7.0, 0.3,
             "Muhajir Project Tilawah · Walkthrough Assessment Al-Fatihah", size=9, color=c)
    add_text(slide, 10.5, 6.95, 2.2, 0.3, f"{_page} / {TOTAL}", size=9, color=c,
             align=PP_ALIGN.RIGHT)


def place_image(slide, name, x, y, max_w, max_h, crop=None, border=True):
    """Sisipkan gambar, skala proporsional agar muat dalam kotak yang diberikan.

    `crop` = (atas, bawah) sebagai pecahan tinggi — untuk tangkapan layar
    panjang yang harus dipenggal jadi beberapa slide.
    """
    path = IMG / name
    if not path.exists():
        add_text(slide, x, y, max_w, 0.4, f"[gambar hilang: {name}]", size=11, color=DANGER)
        return None

    img = Image.open(path)
    src = path
    if crop:
        top_f, bot_f = crop
        w, h = img.size
        img = img.crop((0, int(h * top_f), w, int(h * bot_f)))
        src = ROOT / ".venv" / f"_crop_{name}"
        img.save(src)

    iw, ih = img.size
    scale = min(max_w / (iw / 96), max_h / (ih / 96))
    w_in, h_in = (iw / 96) * scale, (ih / 96) * scale
    cx = x + (max_w - w_in) / 2

    if border:
        add_box(slide, cx - 0.05, y - 0.05, w_in + 0.1, h_in + 0.1, fill=WHITE, border=LINE)
    slide.shapes.add_picture(str(src), Inches(cx), Inches(y), Inches(w_in), Inches(h_in))
    return w_in


def wrapped_height(text, width_in, size_pt, spacing=1.35):
    """Tinggi yang benar-benar dipakai sebuah paragraf setelah membungkus baris.

    python-pptx tidak bisa mengukur teks — tinggi kotak yang kita beri hanya
    saran, teks tetap meluber melewatinya. Jadi jarak antar butir harus dihitung
    sendiri, kalau tidak butir panjang akan menimpa butir di bawahnya.

    Angka 0,55 bukan tebakan: diukur balik dari hasil render — kolom selebar
    3,66 inci memuat 40 huruf pada 11,5 pt. Nilai 0,48 yang biasa dikutip untuk
    Calibri terlalu ramping, dan membuat butir panjang menimpa butir di
    bawahnya. Faktor 1,05 sisanya untuk berjaga.
    """
    char_w = 0.55 * size_pt * 1.05 / 72
    per_line = max(8, int(width_in / char_w))
    baris = max(1, -(-len(text) // per_line))
    return baris * (size_pt * spacing / 72)


def caption(slide, x, y, w, lines):
    """Keterangan di samping gambar: judul tebal + butir-butir."""
    yy = y
    for txt, kind in lines:
        if kind == "h":
            h = wrapped_height(txt, w, 15, 1.2)
            add_text(slide, x, yy, w, h + 0.1, txt, size=15, bold=True, color=PRIMARY,
                     line_spacing=1.2)
            yy += h + 0.16
        elif kind == "route":
            b = add_box(slide, x, yy, min(w, 0.085 * len(txt) + 0.3), 0.28,
                        fill=SURFACE_DEEP, border=None)
            b.adjustments[0] = 0.3
            add_text(slide, x + 0.14, yy, w, 0.28, txt, size=9.5, bold=True,
                     color=INK_SOFT, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
            yy += 0.46
        else:
            h = wrapped_height(txt, w - 0.24, 11.5)
            add_text(slide, x + 0.24, yy, w - 0.24, h + 0.1, txt, size=11.5,
                     color=INK_SOFT, line_spacing=1.35)
            add_circle(slide, x + 0.03, yy + 0.08, 0.09, ACCENT)
            yy += h + 0.17
    return yy


def walkthrough_slide(eyebrow, title, subtitle, image, route, bullets,
                      crop=None, wide=False):
    """Pola slide utama: gambar di kiri, keterangan di kanan."""
    s = new_slide()
    header(s, eyebrow, title, subtitle)
    if wide:
        place_image(s, image, 0.6, 2.0, 7.9, 4.6, crop=crop)
        cap_x, cap_w = 8.9, 3.9
    else:
        place_image(s, image, 0.7, 1.95, 3.5, 4.7, crop=crop)
        cap_x, cap_w = 4.7, 8.0
    y = 2.05
    if route:
        y = caption(s, cap_x, y, cap_w, [(route, "route")])
    caption(s, cap_x, y, cap_w, bullets)
    footer(s)
    return s


# =========================================================
# 1 — Sampul
# =========================================================
s = new_slide(PRIMARY)
add_text(s, 0.9, 1.5, 11.5, 0.6, "بِسْمِ ٱللَّهِ ٱلرَّحْمَٰنِ ٱلرَّحِيمِ", size=26,
         color=ACCENT_SOFT, align=PP_ALIGN.CENTER, font="Traditional Arabic")
add_text(s, 0.9, 2.5, 11.5, 0.4, "WALKTHROUGH SISTEM", size=12, bold=True,
         color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 3.0, 11.5, 1.0, "Assessment Al-Fatihah", size=46, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 4.1, 11.5, 0.4,
         "Perjalanan peserta dan pengajar, dari merekam sampai rapot diterima",
         size=15, color=ACCENT_SOFT, align=PP_ALIGN.CENTER)
bar = add_box(s, 5.4, 4.75, 2.5, 0.04, fill=ACCENT, border=None, radius=False)
add_text(s, 0.9, 5.15, 11.5, 0.35,
         "Sudah berjalan di production · 4 Agustus 2026", size=12,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 5.6, 11.5, 0.35,
         "Tangkapan layar diambil langsung dari sistem yang hidup",
         size=10.5, color=INK_MUTE, align=PP_ALIGN.CENTER)
footer(s, dark=True)

# =========================================================
# 2 — Peta alur
# =========================================================
s = new_slide()
header(s, "Peta", "Satu halaman untuk seluruh alur",
       "Tiga pelaku: peserta, sistem, pengajar. Semua langkah di bawah sudah otomatis kecuali yang bertanda tangan manusia.")

lanes = [
    ("PESERTA", PRIMARY, 1.95),
    ("SISTEM", ACCENT, 3.40),
    ("PENGAJAR", SUCCESS, 4.85),
]
for nama, warna, y in lanes:
    add_box(s, 0.6, y, 1.55, 1.25, fill=warna, border=None)
    add_text(s, 0.6, y, 1.55, 1.25, nama, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

steps = [
    (2.45, 1.95, "1 · Rekam\nbacaan", PRIMARY),
    (4.55, 1.95, "2 · Isi data\n+ nomor WA", PRIMARY),
    (6.65, 3.40, "3 · Audio disimpan\n+ tugas dibagi", ACCENT),
    (8.75, 3.40, "4 · WA ke\ndua pihak", ACCENT),
    (6.65, 4.85, "5 · Dengar +\ncentang temuan", SUCCESS),
    (8.75, 4.85, "6 · Simpan\npenilaian", SUCCESS),
    (10.85, 1.95, "7 · Rapot\nditerima", PRIMARY),
]
for x, y, teks, warna in steps:
    add_box(s, x, y, 1.9, 1.25, fill=WHITE, border=warna, border_w=1.5)
    add_text(s, x + 0.12, y, 1.66, 1.25, teks, size=11, bold=True, color=warna,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

for x, y in [(4.4, 2.45), (10.7, 2.45), (8.6, 3.90), (8.6, 5.35)]:
    add_arrow(s, x, y, 0.3, 0.24)

add_box(s, 0.6, 6.35, 12.1, 0.45, fill=SURFACE_DEEP, border=None)
add_text(s, 0.85, 6.35, 11.6, 0.45,
         "Langkah 5 satu-satunya yang dikerjakan manusia. Sisanya berjalan sendiri.",
         size=11.5, bold=True, color=INK_SOFT, anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# =========================================================
# BAGIAN A — PESERTA
# =========================================================
s = new_slide(PRIMARY)
add_text(s, 0.9, 2.9, 11.5, 0.4, "BAGIAN A", size=13, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)
add_text(s, 0.9, 3.4, 11.5, 0.8, "Perjalanan Peserta", size=40, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 4.4, 11.5, 0.4,
         "Empat layar, sekitar tiga menit — lalu menunggu kabar lewat WhatsApp",
         size=14, color=ACCENT_SOFT, align=PP_ALIGN.CENTER)
footer(s, dark=True)

walkthrough_slide(
    "Langkah 1", "Halaman depan", "Titik masuk. Peserta mengenal apa yang akan dinilai sebelum mulai.",
    "01-landing.png", "/",
    [("Yang dilihat peserta", "h"),
     ("Ajakan utama “Mulai Assessment Gratis” langsung menuju halaman persetujuan.", "b"),
     ("Empat indikator dijelaskan lebih dulu, jadi peserta tahu yang dinilai apa saja.", "b"),
     ("Catatan: sebagian teks di sini masih menjanjikan rapot AI 30 detik — perlu disesuaikan dengan alur baru yang dinilai pengajar.", "b")],
)

walkthrough_slide(
    "Langkah 2", "Persetujuan privasi", "Wajib disetujui sebelum merekam. Ini yang memenuhi UU PDP.",
    "02-consent.png", "/assessment/consent",
    [("Empat janji yang dinyatakan terbuka", "h"),
     ("Audio disimpan paling lama 7 hari, lalu terhapus otomatis.", "b"),
     ("Rekaman tidak dibagikan ke pihak ketiga.", "b"),
     ("Nama dan nomor WA dipakai hanya untuk mengirim hasil.", "b"),
     ("Hasilnya referensi belajar, bukan pengganti pengajar.", "b")],
)

walkthrough_slide(
    "Langkah 3", "Merekam bacaan", "Teks Al-Fatihah ditampilkan utuh supaya peserta tidak perlu membuka mushaf lain.",
    "03-record.png", "/assessment/record",
    [("Yang membantu peserta", "h"),
     ("Tujuh ayat tampil lengkap dengan harakat, plus tombol terjemahan.", "b"),
     ("Boleh direkam ulang sebanyak yang diperlukan — tidak ada batas percobaan.", "b"),
     ("Gelombang suara bergerak saat merekam, jadi peserta yakin suaranya tertangkap.", "b")],
)

# Layar form: tidak bisa difoto tanpa mikrofon nyata — dijelaskan apa adanya
s = new_slide()
header(s, "Langkah 4", "Data peserta",
       "Layar terakhir sebelum rekaman dikirim. Hanya tiga isian.")
add_box(s, 0.7, 2.0, 5.6, 4.3, fill=WHITE, border=LINE)
add_text(s, 1.0, 2.3, 5.0, 0.3, "TIGA ISIAN", size=10, bold=True, color=ACCENT)
fields = [
    ("Nama", "Dipakai menyapa peserta di rapot dan pesan WhatsApp"),
    ("Jenis Kelamin", "Ikhwan atau akhwat — menentukan pengajar yang menilai"),
    ("Nomor WhatsApp", "Awalan +62 sudah terpasang, peserta tinggal melanjutkan"),
]
yy = 2.85
for label, ket in fields:
    add_box(s, 1.0, yy, 5.0, 0.95, fill=SURFACE, border=LINE)
    add_text(s, 1.25, yy + 0.17, 4.5, 0.3, label, size=13, bold=True, color=PRIMARY)
    add_text(s, 1.25, yy + 0.5, 4.5, 0.35, ket, size=10.5, color=INK_SOFT)
    yy += 1.15
caption(s, 6.9, 2.05, 5.8, [
    ("/assessment/form", "route"),
    ("Kenapa layar ini tidak difoto", "h"),
    ("Halaman ini hanya bisa dibuka kalau ada rekaman di memori browser. Rekaman itu sengaja TIDAK disimpan ke penyimpanan browser — keputusan privasi, supaya audio hilang saat tab ditutup.", "b"),
    ("Akibatnya layar ini tidak bisa dijangkau alat otomatis tanpa mikrofon sungguhan. Isinya digambarkan apa adanya di sebelah.", "b"),
    ("Setelah dikirim, peserta diarahkan ke layar tunggu dan langsung menerima WhatsApp.", "b"),
])
footer(s)

walkthrough_slide(
    "Langkah 5", "Konfirmasi masuk WhatsApp",
    "Terkirim otomatis begitu rekaman diterima — peserta tidak dibiarkan menebak.",
    "05-wa-diterima.png", "otomatis · lib/whatsapp.ts",
    [("Kenapa pesan ini ada", "h"),
     ("Penilaian dikerjakan manusia, jadi jedanya berhari-hari, bukan 30 detik. Tanpa kabar, peserta akan mengira sistemnya rusak.", "b"),
     ("Tautan di dalamnya sudah menuju halaman rapot — halaman yang sama nantinya berisi hasil.", "b"),
     ("Tampilan chat direkayasa; kalimatnya persis yang dikirim sistem.", "b")],
    wide=True,
)

# =========================================================
# BAGIAN B — PENGAJAR
# =========================================================
s = new_slide(PRIMARY)
add_text(s, 0.9, 2.9, 11.5, 0.4, "BAGIAN B", size=13, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)
add_text(s, 0.9, 3.4, 11.5, 0.8, "Perjalanan Pengajar", size=40, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 4.4, 11.5, 0.4,
         "Dengar sambil menilai, dalam satu layar — target 5 menit per rekaman",
         size=14, color=ACCENT_SOFT, align=PP_ALIGN.CENTER)
footer(s, dark=True)

# Pembagian tugas
s = new_slide()
header(s, "Di balik layar", "Siapa yang kebagian menilai",
       "Dipilih sistem, bukan diatur manual. Dua aturan, tanpa pengecualian.")
cards = [
    ("Gender ketat", "Rekaman ikhwan hanya ke pengajar ikhwan, akhwat hanya ke akhwat. Kalau tidak ada yang cocok, sistem TIDAK menyeberang gender — tugas jatuh ke superadmin.", SUCCESS),
    ("Giliran merata", "Yang paling lama tidak kebagian dapat giliran duluan. Pengajar yang belum pernah menerima didahulukan.", ACCENT),
    ("Tidak ada yang hilang", "Setiap tugas tercatat di daftar, jadi rekaman tetap terjangkau walau pesan WhatsApp-nya terhapus.", PRIMARY),
]
x = 0.75
for judul, isi, warna in cards:
    add_box(s, x, 2.15, 3.85, 3.1, fill=WHITE, border=LINE)
    add_box(s, x, 2.15, 3.85, 0.09, fill=warna, border=None, radius=False)
    add_text(s, x + 0.35, 2.55, 3.15, 0.4, judul, size=16, bold=True, color=PRIMARY)
    add_text(s, x + 0.35, 3.15, 3.15, 1.9, isi, size=11.5, color=INK_SOFT, line_spacing=1.4)
    x += 4.05
add_box(s, 0.75, 5.6, 11.85, 0.75, fill=SURFACE_DEEP, border=None)
add_text(s, 1.1, 5.6, 11.2, 0.75,
         "Kalau daftar pengajar belum terisi, semua rekaman jatuh ke nomor superadmin — tidak ada rekaman yang menggantung tanpa penilai.",
         size=12, bold=True, color=INK_SOFT, anchor=MSO_ANCHOR.MIDDLE)
footer(s)

walkthrough_slide(
    "Langkah 6", "Pemberitahuan tugas masuk",
    "Pengajar menerima WhatsApp berisi identitas peserta dan tautan langsung ke halaman penilaian.",
    "07a-wa-pengajar.png", "otomatis · rotasi gender-ketat",
    [("Isi pesannya", "h"),
     ("Nama peserta, gender, dan durasi rekaman — cukup untuk tahu apa yang menunggu.", "b"),
     ("Tautannya membuka langsung halaman penilaian rekaman itu, bukan halaman depan portal.", "b"),
     ("Pengajar tetap harus login lebih dulu; tautan saja tidak memberi akses.", "b")],
    wide=True,
)

walkthrough_slide(
    "Langkah 7", "Masuk portal pengajar",
    "Nomor WhatsApp dan password yang diberikan admin. Alamat portalnya sengaja tidak mudah ditebak dan diblokir dari mesin pencari.",
    "07-login-pengajar.png", "/portal-mpt-x7/login",
    [("Catatan", "h"),
     ("Nomor boleh diketik 0812…, 62812…, atau +62812… — semuanya diterima.", "b"),
     ("Lupa password direset manual oleh admin, tidak ada email pemulihan.", "b"),
     ("Sesi berlaku 12 jam.", "b")],
    wide=True,
)

walkthrough_slide(
    "Langkah 8", "Daftar tugas penilaian",
    "Semua rekaman yang menunggu, terlama di atas — karena peserta itulah yang paling lama menanti.",
    "08-daftar-tugas.png", "/portal-mpt-x7/tugas",
    [("Yang ditampilkan", "h"),
     ("Nama peserta, gender, durasi, dan sudah berapa lama menunggu.", "b"),
     ("Tugas yang terlalu lama menganggur ditandai merah.", "b"),
     ("Bagian “Selesai” menyimpan riwayat penilaian yang sudah dikirim.", "b"),
     ("Halaman ini yang membuat rekaman tetap terjangkau kalau pesan WhatsApp hilang.", "b")],
    wide=True,
)

walkthrough_slide(
    "Langkah 9", "Halaman penilaian",
    "Pemutar rekaman menempel di atas layar dan tidak ikut tergulir — inilah alasan formulirnya dibangun sendiri.",
    "09-form-penilaian.png", "/portal-mpt-x7/nilai/[id]",
    [("Rancangan yang menentukan", "h"),
     ("Pengajar mendengarkan sambil mencentang, tanpa berpindah tab atau menggulir balik.", "b"),
     ("Al-Fatihah dibagi delapan segmen — ayat 7 dipecah dua karena panjang dan memuat dua kelompok kesalahan berbeda.", "b"),
     ("Segmen yang belum ada temuannya tertutup, supaya layar tidak jadi dinding teks.", "b")],
    wide=True,
)

walkthrough_slide(
    "Langkah 10", "Mencentang temuan",
    "Satu segmen dibuka. Semua pilihan sudah tersedia — pengajar tinggal klik, tidak ada kolom ketik bebas.",
    "10-segmen-terbuka.png", "110 pilihan · 53 fatal · 57 perlu diperhatikan",
    [("Cara bacanya", "h"),
     ("“Perlu Diperhatikan” diletakkan lebih dulu karena jauh lebih sering ditemui.", "b"),
     ("Tiap pilihan menyebut huruf dan katanya, dengan label aspek di kanan.", "b"),
     ("Skor berjalan di bawah layar dan berubah tiap centang.", "b")],
    # Bagian atas tangkapan layar cuma segmen tertutup; yang perlu dilihat
    # adalah ayat 5 yang terbuka beserta daftar centangnya.
    crop=(0.375, 0.70), wide=True,
)

# Cara skor dihitung
s = new_slide()
header(s, "Langkah 11", "Bagaimana skornya dihitung",
       "Rumus yang sama dengan yang selama ini dipakai pengajar secara manual — tidak dirancang ulang.")
add_box(s, 0.75, 2.05, 5.7, 4.3, fill=WHITE, border=LINE)
add_text(s, 1.05, 2.3, 5.1, 0.3, "SKOR TIAP SEGMEN", size=10, bold=True, color=ACCENT)
rules = [
    ("Lebih dari 5 kesalahan fatal", "1", DANGER),
    ("Ada 1 kesalahan fatal", "2", DANGER),
    ("5 atau lebih perlu diperhatikan", "3", WARN),
    ("Ada 1 perlu diperhatikan", "4", WARN),
    ("Bersih", "5", SUCCESS),
]
yy = 2.75
for teks, nilai, warna in rules:
    add_text(s, 1.05, yy, 4.2, 0.3, teks, size=12, color=INK_SOFT)
    add_text(s, 5.35, yy, 0.8, 0.3, nilai, size=15, bold=True, color=warna, align=PP_ALIGN.RIGHT)
    yy += 0.55
add_box(s, 1.05, 5.62, 5.1, 0.42, fill=SURFACE_DEEP, border=None)
add_text(s, 1.25, 5.62, 4.9, 0.42, "Nilai akhir = segmen TERLEMAH, lalu dikali dua",
         size=11.5, bold=True, color=PRIMARY, anchor=MSO_ANCHOR.MIDDLE)

caption(s, 6.95, 2.05, 5.7, [
    ("Kenapa yang terlemah, bukan rata-rata", "h"),
    ("Rata-rata akan menyamarkan satu ayat yang rusak parah di antara tujuh ayat yang bagus. Padahal justru ayat itulah yang perlu diperbaiki lebih dulu.", "b"),
    ("Satu kesalahan fatal menjatuhkan segmen ke 2", "h"),
    ("Ini disengaja: lahn jaliy mengubah makna bacaan, jadi bobotnya tidak sebanding dengan kesalahan ringan.", "b"),
    ("Kenapa dikali dua", "h"),
    ("Supaya tampil pada skala 1–10 yang sudah dikenal peserta, dan tiap tingkat mendarat tepat di batas predikat.", "b"),
])
footer(s)

# =========================================================
# BAGIAN C — RAPOT
# =========================================================
s = new_slide(PRIMARY)
add_text(s, 0.9, 2.9, 11.5, 0.4, "BAGIAN C", size=13, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)
add_text(s, 0.9, 3.4, 11.5, 0.8, "Rapot Peserta", size=40, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 4.4, 11.5, 0.4,
         "Yang dilihat peserta adalah nilai pengajar — bukan nilai mesin",
         size=14, color=ACCENT_SOFT, align=PP_ALIGN.CENTER)
footer(s, dark=True)

walkthrough_slide(
    "Langkah 12", "Kabar rapot siap",
    "Terkirim otomatis begitu pengajar menekan simpan. Hanya sekali, walau penilaian diralat.",
    "11-wa-rapot-siap.png", "otomatis · setelah penilaian tersimpan",
    [("Yang dijaga di sini", "h"),
     ("Pesan ini tidak terkirim dua kali. Pengajar boleh meralat penilaiannya, tapi peserta tidak akan dikabari ulang.", "b"),
     ("Tautannya sama dengan yang dikirim di awal — peserta tidak perlu menyimpan dua alamat.", "b")],
    wide=True,
)

walkthrough_slide(
    "Langkah 13", "Rapot diterima",
    "Nilai, predikat, dan penjelasan singkat. Angka ini datang dari pengajar yang mendengarkan langsung.",
    "06a-rapot-skor.png", "/rapot/[slug]",
    [("Bagian atas rapot", "h"),
     ("Nilai besar dengan predikatnya, plus nama pengajar yang memeriksa.", "b"),
     ("Kalimat penjelas ditulis menyemangati, bukan menghakimi — apa pun nilainya.", "b"),
     ("Peserta tidak melihat istilah lahn jaliy atau khafiy; yang tampil “Fatal” dan “Perlu Diperhatikan”.", "b")],
)

s = new_slide()
header(s, "Langkah 14", "Isi rapot selengkapnya",
       "Peta bacaan per bagian, rincian tiap ayat, dan rangkuman lima aspek.")
place_image(s, "06-rapot-penuh.png", 0.7, 1.95, 2.7, 4.7, crop=(0.12, 0.30))
place_image(s, "06-rapot-penuh.png", 3.6, 1.95, 2.7, 4.7, crop=(0.30, 0.48))
caption(s, 6.5, 2.05, 6.2, [
    ("Peta bacaan per bagian", "h"),
    ("Grafik jaring menunjukkan delapan segmen sekaligus. Makin jauh titiknya dari pusat, makin baik bacaan di bagian itu.", "b"),
    ("Penunjuk bagian terlemah", "h"),
    ("Rapot menyebut terus terang bagian mana yang menahan nilainya, supaya peserta tahu harus memperbaiki apa lebih dulu — bukan menebak.", "b"),
    ("Rincian tiap ayat", "h"),
    ("Teks Arab, transliterasi, dan daftar temuan dengan label aspeknya. Ayat yang bersih ditandai jelas, tidak dibiarkan kosong tanpa keterangan.", "b"),
    ("Bisa disimpan", "h"),
    ("Tombol simpan PDF mencetak rapot tanpa tombol dan ajakan program.", "b"),
])
footer(s)

# =========================================================
# Penutup
# =========================================================
s = new_slide()
header(s, "Ringkasan", "Yang berjalan sendiri",
       "Tidak ada langkah di bawah ini yang perlu dikerjakan manual.")
items = [
    ("Pembagian tugas", "Rotasi gender-ketat, yang paling lama menganggur dapat giliran duluan."),
    ("Dua notifikasi WhatsApp", "Ke pengajar saat rekaman masuk, ke peserta saat rapot siap."),
    ("Konfirmasi ke peserta", "Terkirim di detik yang sama dengan rekaman diterima."),
    ("Penghapusan audio", "Rekaman terhapus otomatis setelah 7 hari, dijalankan penyedia penyimpanan."),
    ("Riwayat peserta", "Peserta yang dinilai lebih dari sekali otomatis melihat perbandingan nilainya di rapot — supaya kemajuannya terlihat."),
]
yy = 2.1
for judul, isi in items:
    add_box(s, 0.75, yy, 11.85, 0.82, fill=WHITE, border=LINE)
    add_box(s, 0.75, yy, 0.08, 0.82, fill=ACCENT, border=None, radius=False)
    add_text(s, 1.15, yy + 0.13, 3.4, 0.3, judul, size=13, bold=True, color=PRIMARY)
    add_text(s, 4.7, yy + 0.15, 7.6, 0.55, isi, size=11.5, color=INK_SOFT, line_spacing=1.3)
    yy += 0.95
footer(s)

s = new_slide()
header(s, "Status", "Sudah jalan, dan yang belum",
       "Per 4 Agustus 2026.")
add_box(s, 0.75, 2.05, 5.9, 4.3, fill=WHITE, border=LINE)
add_text(s, 1.05, 2.35, 5.3, 0.3, "SUDAH BERJALAN", size=10, bold=True, color=SUCCESS)
sudah = [
    "Rekam, simpan, dan kirim rekaman",
    "Notifikasi WhatsApp dua arah",
    "Pembagian tugas gender-ketat",
    "Formulir penilaian 8 segmen",
    "Rapot peserta + simpan PDF",
    "Lima pengajar sudah terdaftar",
]
yy = 2.85
for t in sudah:
    add_circle(s, 1.05, yy + 0.05, 0.16, SUCCESS)
    add_text(s, 1.4, yy, 4.9, 0.32, t, size=12, color=INK_SOFT)
    yy += 0.52

add_box(s, 7.05, 2.05, 5.55, 4.3, fill=WHITE, border=LINE)
add_text(s, 7.35, 2.35, 4.9, 0.3, "BELUM SELESAI", size=10, bold=True, color=WARN)
belum = [
    "Teks halaman depan masih menjanjikan rapot AI 30 detik",
    "Halaman “Hasil Assessment” di dashboard peserta belum tersambung ke nilai pengajar",
    "Penilaian AI pembanding belum berjalan otomatis di production",
]
yy = 2.85
for t in belum:
    add_circle(s, 7.35, yy + 0.05, 0.16, WARN)
    add_text(s, 7.7, yy, 4.6, 0.9, t, size=12, color=INK_SOFT, line_spacing=1.3)
    yy += 0.52 + 0.28 * (len(t) // 44)
footer(s)

s = new_slide(PRIMARY)
add_text(s, 0.9, 2.7, 11.5, 0.5, "Alhamdulillah", size=30, bold=True,
         color=ACCENT_SOFT, align=PP_ALIGN.CENTER, font="Traditional Arabic")
add_text(s, 0.9, 3.6, 11.5, 0.6,
         "Rekaman peserta kini sampai ke pengajar,",
         size=19, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 4.05, 11.5, 0.6,
         "dan hasilnya kembali ke peserta — tanpa satu pun langkah manual di antaranya.",
         size=19, color=WHITE, align=PP_ALIGN.CENTER)
add_box(s, 5.9, 4.95, 1.5, 0.03, fill=ACCENT, border=None, radius=False)
add_text(s, 0.9, 5.35, 11.5, 0.35,
         "mpt-web-280965384396.asia-southeast2.run.app", size=12,
         color=ACCENT_SOFT, align=PP_ALIGN.CENTER, font="Consolas")
footer(s, dark=True)

# =========================================================
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))

# Bersihkan berkas potongan sementara
for f in (ROOT / ".venv").glob("_crop_*.png"):
    f.unlink()

print(f"Tersimpan: {OUT}")
print(f"Jumlah slide: {len(prs.slides._sldIdLst)}  (nomor halaman terakhir: {_page})")
